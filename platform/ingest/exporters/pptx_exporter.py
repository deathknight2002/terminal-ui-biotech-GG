"""
PowerPoint Exporter

Generates PowerPoint presentations from financial data and valuation results.
"""

from typing import Dict, Any, List, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging

logger = logging.getLogger(__name__)


class PPTXExporter:
    """
    Export financial data to PowerPoint presentations.
    
    Supports multiple presentation templates:
    - Banker deck
    - Board presentation
    - Investor update
    """
    
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize exporter with optional template.
        
        Args:
            template_path: Path to PowerPoint template file (optional)
        """
        if template_path:
            self.presentation = Presentation(template_path)
        else:
            self.presentation = Presentation()
        
        self.template_path = template_path
    
    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """
        Add title slide.
        
        Args:
            title: Main title
            subtitle: Subtitle text
        """
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        if subtitle and len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle
        
        logger.info(f"Added title slide: {title}")
    
    def add_content_slide(self, title: str, content: List[str]) -> None:
        """
        Add content slide with bullet points.
        
        Args:
            title: Slide title
            content: List of bullet points
        """
        slide_layout = self.presentation.slide_layouts[1]  # Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        # Add content
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
        
        logger.info(f"Added content slide: {title}")
    
    def add_valuation_summary_slide(
        self, 
        ticker: str,
        valuation_data: Dict[str, Any]
    ) -> None:
        """
        Add valuation summary slide.
        
        Args:
            ticker: Company ticker
            valuation_data: Valuation results
        """
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{ticker} Valuation Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add valuation metrics
        metrics_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(4)
        )
        metrics_frame = metrics_box.text_frame
        
        dcf_value = valuation_data.get('dcf', {}).get('enterprise_value', 0)
        multiples_value = valuation_data.get('multiples', {}).get('implied_value', 0)
        
        metrics_text = f"""
DCF Valuation: ${dcf_value:,.0f}M
Multiples Valuation: ${multiples_value:,.0f}M

WACC: {valuation_data.get('wacc', 0):.1%}
Terminal Growth Rate: {valuation_data.get('tgr', 0):.1%}

Scenario: {valuation_data.get('scenario', 'Base')}
Run Date: {valuation_data.get('run_date', datetime.now().strftime('%Y-%m-%d'))}
        """
        
        metrics_frame.text = metrics_text.strip()
        for paragraph in metrics_frame.paragraphs:
            paragraph.font.size = Pt(20)
        
        logger.info(f"Added valuation summary slide for {ticker}")
    
    def add_price_targets_slide(self, price_targets: List[Dict[str, Any]]) -> None:
        """
        Add price targets slide.
        
        Args:
            price_targets: List of price target records
        """
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Analyst Price Targets"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add table
        rows = len(price_targets) + 1
        cols = 4
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(1), Inches(1.5),
            Inches(8), Inches(0.4 * rows)
        ).table
        
        # Set column widths
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(2)
        
        # Header row
        headers = ['Source', 'Date', 'Price Target', 'Currency']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Data rows
        for row_idx, target in enumerate(price_targets, start=1):
            table.cell(row_idx, 0).text = str(target.get('source', ''))
            table.cell(row_idx, 1).text = str(target.get('date', ''))
            table.cell(row_idx, 2).text = f"${target.get('price_target', 0):.2f}"
            table.cell(row_idx, 3).text = str(target.get('currency', 'USD'))
            
            # Format data cells
            for col_idx in range(4):
                cell = table.cell(row_idx, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        logger.info(f"Added price targets slide with {len(price_targets)} targets")
    
    def add_loe_timeline_slide(self, loe_events: List[Dict[str, Any]]) -> None:
        """
        Add LoE timeline slide.
        
        Args:
            loe_events: List of LoE event records
        """
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        title_placeholder.text = "Loss of Exclusivity Timeline"
        
        # Add events as bullet points
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for event in loe_events:
            asset_name = event.get('asset_name', '')
            expiry_date = event.get('expiry_date', '')
            region = event.get('region', '')
            
            p = text_frame.add_paragraph()
            p.text = f"{asset_name} ({region}) - {expiry_date}"
            p.level = 0
            p.font.size = Pt(18)
        
        logger.info(f"Added LoE timeline slide with {len(loe_events)} events")
    
    def generate_banker_deck(
        self,
        ticker: str,
        valuation_data: Dict[str, Any],
        price_targets: Optional[List[Dict[str, Any]]] = None,
        loe_events: Optional[List[Dict[str, Any]]] = None
    ) -> None:
        """
        Generate complete banker deck presentation.
        
        Args:
            ticker: Company ticker
            valuation_data: Valuation results
            price_targets: Optional price targets
            loe_events: Optional LoE events
        """
        # Title slide
        self.add_title_slide(
            f"{ticker} Investment Analysis",
            f"Prepared {datetime.now().strftime('%B %d, %Y')}"
        )
        
        # Executive summary
        self.add_content_slide(
            "Executive Summary",
            [
                "Comprehensive valuation analysis",
                "DCF and multiples approach",
                "Street consensus integration",
                "Risk factors and sensitivity"
            ]
        )
        
        # Valuation summary
        self.add_valuation_summary_slide(ticker, valuation_data)
        
        # Price targets if available
        if price_targets:
            self.add_price_targets_slide(price_targets)
        
        # LoE timeline if available
        if loe_events:
            self.add_loe_timeline_slide(loe_events)
        
        # Appendix
        self.add_content_slide(
            "Appendix",
            [
                "Detailed assumptions",
                "Sensitivity tables",
                "Comparable companies",
                "Risk factors"
            ]
        )
        
        logger.info(f"Generated banker deck for {ticker}")
    
    def save(self, pptx_path: str) -> None:
        """
        Save presentation to file.
        
        Args:
            pptx_path: Output PowerPoint file path
        """
        self.presentation.save(pptx_path)
        logger.info(f"Saved presentation to: {pptx_path}")
